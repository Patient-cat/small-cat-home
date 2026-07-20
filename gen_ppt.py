"""Generate the 10-slide fall-risk-focused competition proposal PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────
C_DARK   = RGBColor(0x0D, 0x3B, 0x66)
C_TEAL   = RGBColor(0x1A, 0x93, 0x6F)
C_WARM   = RGBColor(0xF4, 0xA2, 0x61)
C_LIGHT  = RGBColor(0xF7, 0xF9, 0xFC)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_T = RGBColor(0x1E, 0x29, 0x3B)
C_MUTED  = RGBColor(0x64, 0x74, 0x8B)
C_RED    = RGBColor(0xEF, 0x44, 0x44)
C_GREEN  = RGBColor(0x10, 0xB9, 0x81)

FONT_CN = 'Microsoft YaHei'
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BASE = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(BASE, '方案设计PPT-跌倒风险.pptx')

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ────────────────────────────────────────────────
def add_blank_slide(bg_color=C_LIGHT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    return slide

def add_text_box(slide, text, left, top, width, height, font_size=14,
                 color=C_DARK_T, bold=False, align=PP_ALIGN.LEFT,
                 font_name=FONT_CN):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rich_text(slide, runs, left, top, width, height, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = align
    for i, (txt, fs, clr, bld, fn) in enumerate(runs):
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(fs)
        run.font.color.rgb = clr
        run.font.bold = bld
        run.font.name = fn or FONT_CN
    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape

def add_top_bar(slide):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.06, fill_color=C_TEAL)

def add_page_number(slide, num):
    add_text_box(slide, str(num), 12.5, 7.05, 0.6, 0.35, font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════
s1 = add_blank_slide(C_DARK)
add_shape(s1, MSO_SHAPE.RECTANGLE, 0, 2.0, 13.333, 0.08, fill_color=C_TEAL)

add_text_box(s1, '基于多模态AI监测的老年人',
             0.8, 2.5, 11.7, 1.2, font_size=40, color=C_WHITE, bold=True)
add_text_box(s1, '跌倒风险识别及预警研究',
             0.8, 3.4, 11.7, 1.0, font_size=40, color=C_WHITE, bold=True)

add_shape(s1, MSO_SHAPE.RECTANGLE, 0.8, 4.7, 0.08, 0.8, fill_color=C_TEAL)
add_rich_text(s1, [
    ("赛题 XH-202617\n", 18, C_TEAL, True, FONT_CN),
    ("海康威视 · 萤石网络  揭榜挂帅专项赛", 14, C_MUTED, False, FONT_CN),
], 1.2, 4.7, 8, 0.9)

add_text_box(s1, '研究方向：跌倒风险', 0.8, 5.7, 5, 0.5, font_size=15, color=C_WARM, bold=True)

add_rich_text(s1, [
    ("广东理工学院  XXX团队\n", 16, C_WHITE, True, FONT_CN),
    ("2026年6月", 13, C_MUTED, False, FONT_CN),
], 0.8, 6.2, 5, 0.8)

add_shape(s1, MSO_SHAPE.OVAL, 9.5, 1.5, 3.5, 3.5, fill_color=C_TEAL)
add_shape(s1, MSO_SHAPE.OVAL, 11.0, 3.5, 2.0, 2.0, fill_color=C_WARM)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Background
# ══════════════════════════════════════════════════════════════
s2 = add_blank_slide(C_LIGHT)
add_top_bar(s2)
add_text_box(s2, '问题背景', 0.8, 0.4, 4, 0.7, font_size=32, color=C_DARK, bold=True)
add_page_number(s2, 2)

# Quote
add_shape(s2, MSO_SHAPE.RECTANGLE, 0.8, 1.3, 11.7, 0.85, fill_color=C_WHITE)
add_shape(s2, MSO_SHAPE.RECTANGLE, 0.8, 1.3, 0.06, 0.85, fill_color=C_TEAL)
add_text_box(s2, '"满足老年人多方面需求，让老年人能有一个幸福美满的晚年，'
             '是各级党委和政府的重要责任。"\n—— 习近平总书记',
             1.1, 1.35, 11.2, 0.75, font_size=13, color=C_DARK_T)

# Big stat callout
add_shape(s2, MSO_SHAPE.RECTANGLE, 0.8, 2.5, 5.6, 4.0, fill_color=C_WHITE)
add_shape(s2, MSO_SHAPE.RECTANGLE, 0.8, 2.5, 5.6, 0.06, fill_color=C_RED)
add_text_box(s2, '跌倒 — 老年人意外伤害的首要原因', 1.2, 2.8, 5, 0.45, font_size=18, color=C_RED, bold=True)

stats = [
    ('4000万', '每年约4000万老年人\n经历跌倒事件'),
    ('36.84%', '60岁以上老人跌倒后\n骨折比例'),
    ('近4成', '跌倒事件发生在\n家庭环境中'),
    ('事后为主', '现有技术以"事后检测"\n为主，缺少事前预判'),
]
for i, (num, desc) in enumerate(stats):
    y = 3.5 + i * 0.72
    add_text_box(s2, num, 1.2, y, 1.8, 0.5, font_size=22, color=C_RED, bold=True)
    add_text_box(s2, desc, 3.2, y, 3, 0.55, font_size=11, color=C_DARK_T)

# Right side: pain point diagram
add_shape(s2, MSO_SHAPE.RECTANGLE, 7.0, 2.5, 5.5, 4.0, fill_color=C_WHITE)
add_shape(s2, MSO_SHAPE.RECTANGLE, 7.0, 2.5, 5.5, 0.06, fill_color=C_TEAL)
add_text_box(s2, '当前方案的局限', 7.4, 2.8, 4.5, 0.45, font_size=18, color=C_TEAL, bold=True)

limits = [
    '事后检测为主 — 人已倒地才报警',
    '缺乏步态等前置风险指标',
    '单模态（仅RGB视觉），场景覆盖不足',
    '卫生间等隐私区域无法部署摄像头',
    '缺少个人基线，无法感知"异常变化"',
]
for i, item in enumerate(limits):
    add_text_box(s2, f'✗  {item}', 7.4, 3.5 + i * 0.55, 5, 0.4, font_size=12, color=C_DARK_T)

add_shape(s2, MSO_SHAPE.RECTANGLE, 0.8, 6.75, 11.7, 0.55, fill_color=RGBColor(0xFE, 0xF3, 0xC7))
add_text_box(s2, '核心目标：从"跌倒后识别"升级为"以跌倒风险前置防控为核心，融合跌倒前预判、跌倒时预警、跌倒后快速响应"的全流程技术方案',
             1.0, 6.78, 11.3, 0.5, font_size=12, color=RGBColor(0x92, 0x4D, 0x0E), bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Technical Architecture
# ══════════════════════════════════════════════════════════════
s3 = add_blank_slide(C_WHITE)
add_top_bar(s3)
add_text_box(s3, '技术路线', 0.8, 0.4, 5, 0.7, font_size=32, color=C_DARK, bold=True)
add_text_box(s3, '从"事后检测"到"预判→检测→响应"全链路闭环', 0.8, 0.95, 8, 0.4, font_size=13, color=C_MUTED)
add_page_number(s3, 3)

# Platform bar
add_shape(s3, MSO_SHAPE.RECTANGLE, 0.8, 1.5, 11.7, 0.55, fill_color=C_DARK)
add_text_box(s3, '萤石开放平台  设备接入  |  API 服务  |  云存储  |  告警推送  |  AI 算法订阅',
             1.0, 1.52, 11.3, 0.5, font_size=13, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# Three phases
phases_arch = [
    ('事前预判\n（风险感知）', [
        '步态参数提取（步速/步幅/晃动）',
        '7天个人基线自动建立',
        '日常活动量趋势监测',
        '偏离基线 > 2σ → 风险预警',
    ], C_TEAL, C_LIGHT),
    ('事中检测\n（实时监测）', [
        'YOLOv8-pose 17关键点姿态估计',
        '多特征融合跌倒评分（5维加权）',
        '连续帧确认 + 冷却期防抖',
        '分级告警：黄色预警 + 红色确认',
    ], C_WARM, C_WHITE),
    ('事后响应\n（快速处置）', [
        '跌倒事件自动截图 + 数据库记录',
        'WebSocket / SSE 实时推送',
        'AI 事件分析（DeepSeek）',
        '事件溯源 + 视频回放',
    ], C_RED, C_LIGHT),
]
for i, (phase, items, accent, bg) in enumerate(phases_arch):
    x = 0.8 + i * 4.1
    add_shape(s3, MSO_SHAPE.RECTANGLE, x, 2.3, 3.8, 3.5, fill_color=bg)
    add_shape(s3, MSO_SHAPE.RECTANGLE, x, 2.3, 3.8, 0.06, fill_color=accent)
    add_text_box(s3, phase, x + 0.3, 2.55, 3.2, 0.7, font_size=18, color=accent, bold=True)

    for j, item in enumerate(items):
        add_shape(s3, MSO_SHAPE.RECTANGLE, x + 0.3, 3.4 + j * 0.52, 0.18, 0.18, fill_color=accent)
        add_text_box(s3, item, x + 0.6, 3.35 + j * 0.52, 3, 0.42, font_size=11, color=C_DARK_T)

    if i < 2:
        add_text_box(s3, '▸', x + 3.9, 3.7, 0.4, 0.5, font_size=22, color=C_MUTED, align=PP_ALIGN.CENTER)

# Core algorithms callout
add_shape(s3, MSO_SHAPE.RECTANGLE, 0.8, 6.1, 11.7, 0.75, fill_color=C_DARK)
add_text_box(s3, '核心算法：YOLOv8-pose 姿态估计  +  IoU多目标追踪  +  5维特征融合跌倒评分  +  个人基线步态建模  +  InsightFace人脸识别',
             1.0, 6.13, 11.3, 0.6, font_size=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Fall Risk Pre-Prediction (Core)
# ══════════════════════════════════════════════════════════════
s4 = add_blank_slide(C_LIGHT)
add_top_bar(s4)
add_text_box(s4, '跌倒风险前置预判', 0.8, 0.4, 5, 0.7, font_size=32, color=C_DARK, bold=True)
add_text_box(s4, '从"跌倒后识别"向"预判 → 检测 → 响应"全链路升级', 0.8, 0.95, 8, 0.4, font_size=13, color=C_MUTED)
add_page_number(s4, 4)

# Left: existing capabilities
add_shape(s4, MSO_SHAPE.RECTANGLE, 0.8, 1.5, 5.7, 5.3, fill_color=C_WHITE)
add_shape(s4, MSO_SHAPE.RECTANGLE, 0.8, 1.5, 5.7, 0.06, fill_color=C_TEAL)
add_text_box(s4, '现有基础（已开发完成）', 1.2, 1.7, 4.5, 0.4, font_size=18, color=C_TEAL, bold=True)

existing = [
    ('YOLOv8-pose 姿态估计', '17关键点实时检测，CUDA加速推理，每3帧运行一次'),
    ('多特征融合跌倒评分', '躯干倾斜角(35%) + 垂直速度(25%) + 宽高比(20%)\n+ 角加速度(12%) + 头脚Y向位移(8%) → 加权Sigmoid'),
    ('IoU多目标追踪', '支持多人场景，最大丢失30帧，IoU匹配阈值0.3'),
    ('双级告警机制', '黄色预警(P≥0.50) + 红色确认(P≥0.60, 连续2帧)\n+ 5秒冷却期防止重复触发'),
    ('InsightFace人脸识别', 'CLAHE预处理 + 余弦相似度匹配 + 自动学习入库\n+ 多人场景身份关联'),
    ('AI事件分析 & 数据管理', 'DeepSeek/Gemini事件分析 + SQLite事件记录\n+ 7天自动清理 + 永久标记保护'),
]
for i, (label, desc) in enumerate(existing):
    y = 2.3 + i * 0.72
    add_shape(s4, MSO_SHAPE.RECTANGLE, 1.2, y, 0.2, 0.2, fill_color=C_GREEN)
    add_text_box(s4, label, 1.6, y - 0.05, 4.5, 0.3, font_size=12, color=C_DARK_T, bold=True)
    add_text_box(s4, desc, 1.6, y + 0.22, 4.7, 0.4, font_size=10, color=C_MUTED)

# Right: new additions
add_shape(s4, MSO_SHAPE.RECTANGLE, 7.0, 1.5, 5.5, 5.3, fill_color=C_WHITE)
add_shape(s4, MSO_SHAPE.RECTANGLE, 7.0, 1.5, 5.5, 0.06, fill_color=C_WARM)
add_text_box(s4, '新增：步态风险预判模块', 7.4, 1.7, 4.5, 0.4, font_size=18, color=C_WARM, bold=True)

add_text_box(s4, '步态参数提取', 7.4, 2.3, 4.5, 0.3, font_size=14, color=C_DARK_T, bold=True)
gait_params = ['步速（髋关节位移 / 秒）', '步幅（踝关节跨步距离）', '躯干晃动方差', '双足支撑相占比', '起坐耗时']
for i, p in enumerate(gait_params):
    add_text_box(s4, f'▸ {p}', 7.6, 2.65 + i * 0.35, 4.5, 0.3, font_size=11, color=C_DARK_T)

add_text_box(s4, '个人基线建模', 7.4, 4.55, 4.5, 0.3, font_size=14, color=C_DARK_T, bold=True)
baseline_steps = [
    '1. 采集前7天步态参数，建立个人均值 ± 标准差基线',
    '2. 每日自动采样，与基线对比',
    '3. 偏离 > 2σ → 推送"跌倒风险升高"预警',
    '4. 基线持续更新，自适应老人身体变化',
]
for i, step in enumerate(baseline_steps):
    add_text_box(s4, step, 7.6, 4.9 + i * 0.32, 4.7, 0.3, font_size=11, color=C_DARK_T)

add_text_box(s4, '硬件：C6C Ultra ×2（客厅 + 卧室）', 7.4, 6.2, 5, 0.35, font_size=12, color=C_TEAL, bold=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Device Deployment
# ══════════════════════════════════════════════════════════════
s5 = add_blank_slide(C_LIGHT)
add_top_bar(s5)
add_text_box(s5, '硬件部署方案', 0.8, 0.4, 5, 0.7, font_size=32, color=C_DARK, bold=True)
add_text_box(s5, '覆盖居家跌倒高风险区域：视觉区 + 隐私区', 0.8, 0.95, 8, 0.4, font_size=13, color=C_MUTED)
add_page_number(s5, 5)

table_data = [
    ['序号', '设备型号', '数量', '部署位置', '核心用途'],
    ['1', '萤石 C6C Ultra（旗舰AI摄像机）', '2', '客厅、卧室', '姿态估计 + 步态分析 + 跌倒检测'],
    ['2', '萤石 C6C 5MP（高清室内云台）', '1', '走廊', '行为追踪 + 辅助监控'],
    ['3', '萤石 CS-A3-V200-WBG（智能网关）', '1', '全屋中枢', 'Zigbee 3.0 设备联网 + 场景联动'],
    ['4', '萤石 DS-TDSB00-EKT（跌倒检测雷达）', '1', '卫生间', '毫米波跌倒检测 — 零图像，隐私安全'],
    ['5', '萤石 DS-TDSB00-EKH（生命体征雷达）', '1', '卧室', '睡眠/呼吸/心率监测 — 辅助步态评估'],
    ['6', '萤石 T21（人体存在传感器）', '1', '卫生间/厨房', '存在检测 + 停留时长异常告警'],
    ['7', '萤石 T51C（温湿度传感器）', '1', '浴室', '湿度异常 → 地面湿滑风险预警'],
    ['8', '萤石 SP10（智能中控屏）', '1', '客厅', '语音交互 + 一键呼叫 + 告警联动'],
]
rows, cols = len(table_data), len(table_data[0])
tbl = s5.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(12.1), Inches(4.5)).table
tbl.columns[0].width = Inches(0.7)
tbl.columns[1].width = Inches(4.0)
tbl.columns[2].width = Inches(0.8)
tbl.columns[3].width = Inches(2.0)
tbl.columns[4].width = Inches(4.6)

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = FONT_CN
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = C_WHITE
            else:
                p.font.color.rgb = C_DARK_T
            p.alignment = PP_ALIGN.CENTER if c in [0, 2] else PP_ALIGN.LEFT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_DARK
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_LIGHT

# Floor plan
add_shape(s5, MSO_SHAPE.RECTANGLE, 0.8, 6.2, 11.7, 0.7, fill_color=C_WHITE)
add_shape(s5, MSO_SHAPE.RECTANGLE, 0.8, 6.2, 0.06, 0.7, fill_color=C_TEAL)
add_text_box(s5, '部署示意：客厅(C6C Ultra + SP10) → 卧室(C6C Ultra + DS-TDSB00-EKH) → 卫生间(DS-TDSB00-EKT + T21) → 走廊(C6C 5MP) → 浴室(T51C)',
             1.2, 6.25, 11, 0.55, font_size=13, color=C_DARK_T)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — EZVIZ Platform
# ══════════════════════════════════════════════════════════════
s6 = add_blank_slide(C_WHITE)
add_top_bar(s6)
add_text_box(s6, '基于萤石开放平台的实现', 0.8, 0.4, 5, 0.7, font_size=32, color=C_DARK, bold=True)
add_page_number(s6, 6)

layers = [
    ('设备层', '萤石 SDK + 硬件接入', [
        'C6C Ultra / C6C 5MP 视频流拉取',
        'DS-TDSB00 毫米波雷达数据读取',
        'T21 / T51C Zigbee 传感器接入 (CS-A3网关)',
        'SP10 智能屏语音交互通道',
    ], C_TEAL),
    ('算法层', '自研 AI 引擎', [
        'YOLOv8-pose 17关键点姿态估计',
        '5维特征融合跌倒评分算法',
        '步态参数提取 + 个人基线建模',
        'InsightFace 人脸识别 + 自动学习',
    ], C_WARM),
    ('联动层', '萤石 API + 云存储', [
        'WebSocket / SSE 实时告警推送',
        '萤石云存储：事件截图 + 录像',
        '萤石 AI 跌倒检测作为对照基线',
        '设备联动：传感器触发 → 摄像头抓拍 → 告警',
    ], RGBColor(0x63, 0x66, 0xF1)),
    ('展示层', 'Web 管理平台', [
        '监控大厅：多路实时画面 + 姿态骨架叠加',
        '风险看板：步态趋势 + 个人基线偏离',
        '事件溯源：跌倒前后30秒视频回放',
        '分级告警：蓝(异常) → 黄(预警) → 红(确认)',
    ], RGBColor(0xEC, 0x48, 0x99)),
]
for i, (name, sub, items, accent) in enumerate(layers):
    y = 1.3 + i * 1.45
    add_shape(s6, MSO_SHAPE.RECTANGLE, 0.8, y, 11.7, 1.2, fill_color=C_LIGHT)
    add_shape(s6, MSO_SHAPE.RECTANGLE, 0.8, y, 0.06, 1.2, fill_color=accent)
    add_text_box(s6, name, 1.2, y + 0.08, 1.5, 0.35, font_size=16, color=accent, bold=True)
    add_text_box(s6, sub, 2.8, y + 0.08, 3.5, 0.35, font_size=12, color=C_MUTED)
    for j, item in enumerate(items):
        add_text_box(s6, f'▸ {item}', 1.2 + (j % 2) * 5.8, y + 0.5 + (j // 2) * 0.33, 5.3, 0.3, font_size=11, color=C_DARK_T)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Evaluation
# ══════════════════════════════════════════════════════════════
s7 = add_blank_slide(C_LIGHT)
add_top_bar(s7)
add_text_box(s7, '评价指标与验证方案', 0.8, 0.4, 5, 0.7, font_size=32, color=C_DARK, bold=True)
add_page_number(s7, 7)

eval_data = [
    ['指标', '目标值', '验证方式'],
    ['跌倒检测准确率', '≥ 90%', 'CAUCAFall 公开数据集 200+ 段跌倒视频测试'],
    ['跌倒检测召回率', '≥ 85%', '自建居家场景多角度测试（客厅/卧室/走廊）'],
    ['步态风险预判准确率', '≥ 80%', '7天连续监测，与 Berg 平衡量表评分对比验证'],
    ['误报率（日常活动）', '≤ 10%', '24h 连续运行，统计弯腰/躺下/坐下等类跌倒误报次数'],
    ['卫生间跌倒检出率', '≥ 85%', '毫米波雷达 (DS-TDSB00-EKT) 独立测试，模拟跌倒 ×50次'],
    ['系统响应延迟', '≤ 2秒', '从跌倒发生到告警推送的端到端延迟测量'],
]
e_rows, e_cols = len(eval_data), len(eval_data[0])
etbl = s7.shapes.add_table(e_rows, e_cols, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.2)).table
etbl.columns[0].width = Inches(3.2)
etbl.columns[1].width = Inches(1.5)
etbl.columns[2].width = Inches(7.0)

for r in range(e_rows):
    for c in range(e_cols):
        cell = etbl.cell(r, c)
        cell.text = eval_data[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14) if r == 0 else Pt(13)
            p.font.name = FONT_CN
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = C_WHITE
                p.alignment = PP_ALIGN.CENTER
            else:
                p.font.color.rgb = C_DARK_T
                p.alignment = PP_ALIGN.CENTER if c == 1 else PP_ALIGN.LEFT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_DARK
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE

add_text_box(s7, '全部测试将在萤石设备实景部署环境中完成，提供完整实测数据、截图及可复现测试脚本。',
             0.8, 6.0, 11.7, 0.5, font_size=12, color=C_MUTED)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Elderly-Friendly Design
# ══════════════════════════════════════════════════════════════
s8 = add_blank_slide(C_WHITE)
add_top_bar(s8)
add_text_box(s8, '适老化设计', 0.8, 0.4, 5, 0.7, font_size=32, color=C_DARK, bold=True)
add_page_number(s8, 8)

designs = [
    ('无感采集', '不戴手环 · 不填问卷 · 零操作\n所有数据自动采集、自动分析\n老人完全无感知', C_GREEN),
    ('隐私保护', '卫生间仅用毫米波雷达\n(DS-TDSB00-EKT)\n零图像 · 仅点云坐标 · 不识别身份', C_TEAL),
    ('告警分层', '家属 → 社区网格员 → 120\n逐级递进 · 不误报 · 不遗漏\n每级间隔可配置', C_WARM),
    ('一键呼叫', 'SP10 智能屏语音直达家属\n"帮我叫一下儿子" 即刻响应\n支持紧急联系人轮询', RGBColor(0x63, 0x66, 0xF1)),
    ('数据本地化', '面部数据仅存本地 · 不上传云\n行为特征脱敏处理\n符合个人信息保护法要求', RGBColor(0xEC, 0x48, 0x99)),
]
for i, (title, desc, accent) in enumerate(designs):
    x = 0.6 + i * 2.5
    add_shape(s8, MSO_SHAPE.RECTANGLE, x, 1.6, 2.25, 3.5, fill_color=C_LIGHT)
    add_shape(s8, MSO_SHAPE.RECTANGLE, x, 1.6, 2.25, 0.06, fill_color=accent)
    add_text_box(s8, title, x + 0.15, 1.85, 1.95, 0.4, font_size=18, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(s8, desc, x + 0.15, 2.5, 1.95, 2.2, font_size=11, color=C_DARK_T, align=PP_ALIGN.CENTER)

add_shape(s8, MSO_SHAPE.RECTANGLE, 0.8, 5.5, 11.7, 0.55, fill_color=C_DARK)
add_text_box(s8, '设计原则：老人无需学习任何操作，系统自动感知、自动分析、自动预警 — "让技术隐形，让安全可见"',
             1.0, 5.53, 11.3, 0.5, font_size=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Timeline
# ══════════════════════════════════════════════════════════════
s9 = add_blank_slide(C_LIGHT)
add_top_bar(s9)
add_text_box(s9, '团队分工与时间规划', 0.8, 0.4, 5, 0.7, font_size=32, color=C_DARK, bold=True)
add_page_number(s9, 9)

phases = [
    ('5-6月', '方案设计 & 设备对接', [
        '方案设计 + PPT + 资源申请',
        '萤石设备到货 + 环境搭建',
        '开放平台 API 集成联调',
        'C6C Ultra 视频流 + CS-A3 传感器接入',
    ], C_TEAL),
    ('7月', '步态预判 & 算法优化', [
        '步态参数提取模块开发',
        '个人基线建模算法实现',
        '跌倒检测 + 步态预判联合调优',
        '单元测试 + 模块验证',
    ], C_WARM),
    ('8月', '系统联调 & 实景测试', [
        '全设备联调：摄像头 + 雷达 + 传感器',
        '居家实景部署（客厅/卧室/卫生间）',
        '24h 连续运行稳定性测试',
        '误报优化 + 性能调优',
    ], RGBColor(0x63, 0x66, 0xF1)),
    ('9月上旬', '材料提交', [
        '研究报告撰写',
        '源代码整理 + 部署文档',
        '测试报告 + 实测数据',
        '9月5日前完成全部提交',
    ], RGBColor(0xEC, 0x48, 0x99)),
    ('10-11月', '终审准备', [
        '入围后按专家意见优化',
        '终审擂台赛答辩准备',
    ], RGBColor(0xE8, 0x6A, 0x17)),
]
for i, (time, title, items, accent) in enumerate(phases):
    y = 1.3 + i * 1.15
    add_shape(s9, MSO_SHAPE.OVAL, 1.15, y + 0.12, 0.2, 0.2, fill_color=accent)
    if i < len(phases) - 1:
        add_shape(s9, MSO_SHAPE.RECTANGLE, 1.22, y + 0.35, 0.06, 0.85, fill_color=RGBColor(0xE2, 0xE8, 0xF0))
    add_text_box(s9, time, 1.6, y - 0.05, 1.5, 0.35, font_size=14, color=accent, bold=True)
    add_text_box(s9, title, 3.2, y - 0.05, 4, 0.35, font_size=16, color=C_DARK_T, bold=True)
    for j, item in enumerate(items):
        add_text_box(s9, f'▸ {item}', 3.2, y + 0.35 + j * 0.25, 9, 0.25, font_size=11, color=C_DARK_T)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Innovation Summary
# ══════════════════════════════════════════════════════════════
s10 = add_blank_slide(C_DARK)
add_text_box(s10, '创新点总结', 0.8, 0.6, 5, 0.8, font_size=36, color=C_WHITE, bold=True)
add_page_number(s10, 10)

innovations = [
    ('01', '从"事后识别"到"事前预判"',
     '步态参数提取 + 7天个人基线建模，首次实现居家跌倒风险前置预判。\n当步速、步幅、晃动等指标偏离基线超过2σ时，在跌倒发生前推送风险预警，\n构建"预判 → 检测 → 响应"全链路闭环。',
     C_TEAL),
    ('02', '从"单点视觉"到"多模态融合"',
     'RGB摄像机(C6C Ultra) + 毫米波雷达(DS-TDSB00-EKT/EKH) + Zigbee传感器(T21/T51C)\n的晚期融合架构。视觉负责客厅/卧室的姿态分析，雷达负责卫生间的隐私安全检测，\n传感器负责环境风险(湿度/停留时长)，三路信号互补覆盖。',
     C_WARM),
    ('03', '从"统一阈值"到"个人基线"',
     '传统方案对所有老人使用相同跌倒判定阈值。本方案为每位老人建立独立的步态基线，\n能够感知"这个老人今天比他自己平时走得慢/晃得厉害"的个体化异常，\n实现真正的个性化风险预警，大幅降低误报率。',
     RGBColor(0x63, 0x66, 0xF1)),
    ('04', '隐私安全 · 无感设计',
     '卫生间零摄像头 — 仅用毫米波雷达(DS-TDSB00-EKT)输出点云坐标。\n面部数据本地存储不上云，行为特征脱敏处理。老人无需任何操作，\n系统全自动运行 — "让技术隐形，让安全可见"。',
     RGBColor(0xEC, 0x48, 0x99)),
]
for i, (num, title, desc, accent) in enumerate(innovations):
    y = 1.8 + i * 1.35
    add_shape(s10, MSO_SHAPE.RECTANGLE, 0.8, y, 11.7, 1.1, fill_color=RGBColor(0x15, 0x3E, 0x6B))
    add_shape(s10, MSO_SHAPE.RECTANGLE, 0.8, y, 0.08, 1.1, fill_color=accent)
    add_text_box(s10, num, 1.2, y + 0.15, 0.8, 0.8, font_size=30, color=accent, bold=True)
    add_text_box(s10, title, 2.1, y + 0.08, 5, 0.4, font_size=18, color=C_WHITE, bold=True)
    add_text_box(s10, desc, 2.1, y + 0.5, 10, 0.55, font_size=10, color=RGBColor(0x94, 0xA3, 0xB8))

add_text_box(s10, '研究方向：跌倒风险  |  赛题 XH-202617  |  基于萤石开放平台  |  广东理工学院',
             0.8, 6.9, 11.7, 0.35, font_size=12, color=C_MUTED, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────
prs.save(OUTPUT)
print(f'PPT saved: {OUTPUT}')
print(f'Slides: {len(prs.slides)}')
